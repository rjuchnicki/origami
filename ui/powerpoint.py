from pptx import Presentation

prs = Presentation()
title_slidelayout = prs.slidelayouts[0]
slide = prs.slides.add_slide(title_slidelayout)
title = slide.shapes.title
subtitle = slide.shapes.placeholders[1]

title.text = "Hello, World!"
subtitle.text = "python-pptx was here!"

prs.save('test.ppsx')