from pptx import Presentation
from pptx.util import Inches, Pt
from glob import glob

from pptx import Presentation

prs = Presentation('template.pptx')
title_slidelayout = prs.slidelayouts[6]

for filename in glob('*.txt'):
	slide = prs.slides.add_slide(title_slidelayout)
	title = slide.shapes.title
	left = top = width = height = Inches(1)
	txBox = slide.shapes.add_textbox(left, top, width, height)
	tf = txBox.textframe

	p = tf.add_paragraph()
	p.text = open(filename).read().replace('\t', '       ')
	p.font.bold = False

prs.save("/Users/tommychen/Desktop/bah.pptx")
'''
command2 = 'open /Users/tommychen/Desktop/bah.pptx'

from commands import getstatusoutput 
getstatusoutput(command2)
'''