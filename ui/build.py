from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.oxml import _SubElement
from numpy import mean
from commands import getstatusoutput

def execute(command):
	return getstatusoutput(command)

slideshow = Presentation("template.pptx")

def setfont(font, rgb, typeface=None):
    feel = font._Font__rPr
    fill = _SubElement(feel, 'a:solidFill')
    color = _SubElement(fill, 'a:srgbClr')
    color.set('val', rgb)
    if typeface:
        latin = _SubElement(feel, 'a:latin')
        latin.set('typeface', typeface)

def addcode(code, previous):
	blank = slideshow.slidelayouts[6]
	slide = slideshow.slides.add_slide(blank)
	box = slide.shapes.add_textbox(Inches(1), Inches(0.5), Inches(8), Inches(10))
	frame = box.textframe
	snippet = frame.add_paragraph()
	code = code.replace('\t', '     ')
	lines = code.split('\n')
	hotpink = (previous != '' and 
		any(line.strip() in previous for line in lines if line.strip() != ''))
	for line in lines:
		line = line + '\n'
		run = snippet.add_run()	
		run.text = line
		run.font.bold = False
		run.font.size = Pt(25)
		if hotpink and line.strip() not in previous:
			setfont(run.font, 'FF0080', 'Calibri')
		else:
			setfont(run.font, 'FFFFFF', 'Calibri')

def addimg(location):
	blank = slideshow.slidelayouts[6]
	slide = slideshow.slides.add_slide(blank)
	pic = slide.shapes.add_picture(location, Inches(3), Inches(0.5))

def addtitle(tag):
	titles = tag.split('-')
	layout = slideshow.slidelayouts[0]
	slide = slideshow.slides.add_slide(layout)
	slide.shapes.title.text = titles[0]
	slide.shapes[0].textframe.paragraphs[0].runs[0].font.size = Pt(115)
	setfont(slide.shapes[0].textframe.paragraphs[0].runs[0].font, 'FF0080', 'Bariol')
	subtitle = slide.shapes.placeholders[1]
	subtitle.text = titles[1]
	slide.shapes[1].textframe.paragraphs[0].runs[0].font.size = Pt(25)
	setfont(slide.shapes[1].textframe.paragraphs[0].runs[0].font, 'FFFFFF', 'Bariol')

def addbullets(data, previous):
	data = data.replace('\t', '     ')
	lines_all = data.split('\n')
	title = lines_all[0]
	for i in range(2,len(lines_all)+1):
		layout = slideshow.slidelayouts[1]
		slide = slideshow.slides.add_slide(layout)
		slide.shapes.title.text = title
		frame = slide.shapes.placeholders[1].textframe
		lines = lines_all[1:i]
		hotpink = (previous != '' and 
			any(line.strip() in previous for line in lines if line.strip() != ''))
		for line in lines[:-1]:
			p = frame.add_paragraph()
			p.text = line.lstrip(' -')
			p.level = 1
			p.font.bold = False
			p.font.size = Pt(25)
			setfont(p.font, 'FFFFFF', 'Calibri')
		line = lines[-1].lstrip(' -')
		p = frame.add_paragraph()
		p.text = line
		p.level = 1
		p.font.bold = False
		p.font.size = Pt(25)
		setfont(p.font, 'FF0080', 'Calibri')

data = open('slideshow.txt', 'r+')

def next():
	data_here = []
	for line in data:
		if (len(line.strip()) == 0 or 
			not all(c == '-' for c in line.strip())):
			data_here.append(line)
		else:
			break
	if data_here != []:
		combined = ''.join(data_here).strip()
		if combined.strip() == '':
			return next()
		return combined
	return None

previous = ''
nextslide = next()
while nextslide != None:
	try:
		if nextslide.endswith('jpg'):
			addimg(nextslide)
		elif '\n' not in nextslide:
			addtitle(nextslide)
		elif nextslide.split('\n')[1].startswith('-'):
			addbullets(nextslide, previous)
		else:
			addcode(nextslide, previous)
		previous += nextslide
	except:
		pass
	nextslide = next()

slideshow.save('slideshow.pptx')
#execute('open slideshow.pptx')
